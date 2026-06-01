"""파우더 벌크 제조량 예측 프로그램 종합 안내 PPT 생성
바탕화면의 두 텍스트 파일(참고 사항 + 관련 질문)을 종합해 슬라이드로 변환"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ===== 컬러 팔레트 (COSMAX 브랜드) =====
RED = RGBColor(0xc8, 0x10, 0x2e)        # COSMAX 메인
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xf2, 0xf2, 0xf6)
BLUE = RGBColor(0x2d, 0x5b, 0xa8)
GREEN = RGBColor(0x2e, 0x8b, 0x57)
ORANGE = RGBColor(0xf7, 0x7f, 0x00)
WHITE = RGBColor(0xff, 0xff, 0xff)
CODE_BG = RGBColor(0x1e, 0x1e, 0x2e)

# ===== 프레젠테이션 =====
prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=14, color=DARK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='맑은 고딕'):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    return tf


def add_header_bar(slide, title, subtitle=None, accent=RED):
    add_rect(slide, 0, 0, SW, Inches(0.7), accent)
    add_rect(slide, 0, Inches(0.7), SW, Inches(0.04), DARK)
    add_text(slide, Inches(0.5), Inches(0.1), Inches(10), Inches(0.5),
             title, size=22, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(11.0), Inches(0.18), Inches(2.0), Inches(0.4),
                 subtitle, size=11, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_bullet_block(slide, x, y, w, h, items, size=14, accent=RED):
    """items: list of (label, desc) 또는 단일 문자열"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(4)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        if isinstance(item, tuple):
            label, desc = item
            r1 = p.add_run()
            r1.text = '● '
            r1.font.size = Pt(size)
            r1.font.color.rgb = accent
            r1.font.bold = True
            r1.font.name = '맑은 고딕'
            r2 = p.add_run()
            r2.text = label + '  '
            r2.font.size = Pt(size)
            r2.font.color.rgb = DARK
            r2.font.bold = True
            r2.font.name = '맑은 고딕'
            r3 = p.add_run()
            r3.text = desc
            r3.font.size = Pt(size - 1)
            r3.font.color.rgb = GRAY
            r3.font.name = '맑은 고딕'
        else:
            r1 = p.add_run()
            r1.text = '● '
            r1.font.size = Pt(size)
            r1.font.color.rgb = accent
            r1.font.bold = True
            r1.font.name = '맑은 고딕'
            r2 = p.add_run()
            r2.text = item
            r2.font.size = Pt(size)
            r2.font.color.rgb = DARK
            r2.font.name = '맑은 고딕'


def add_code_block(slide, x, y, w, h, code, size=11):
    add_rect(slide, x, y, w, h, CODE_BG)
    box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                   w - Inches(0.3), h - Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = True
    lines = code.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line if line else ' '
        run.font.size = Pt(size)
        run.font.name = 'Consolas'
        run.font.color.rgb = RGBColor(0xe0, 0xe0, 0xe8)


def section_divider(part_label, part_title, part_desc, accent=RED):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, accent)
    add_rect(s, Inches(0.6), Inches(2.6), Inches(0.7), Inches(2.3), WHITE)
    add_text(s, Inches(1.5), Inches(2.5), Inches(2), Inches(0.5),
             part_label, size=14, color=WHITE, bold=True)
    add_text(s, Inches(1.5), Inches(3.0), Inches(11), Inches(1.2),
             part_title, size=44, color=WHITE, bold=True)
    add_text(s, Inches(1.5), Inches(4.5), Inches(11), Inches(1),
             part_desc, size=18, color=WHITE)


# ============================================================
# Slide 1 — 표지
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
add_rect(s, 0, 0, Inches(0.3), SH, RED)
add_rect(s, Inches(0.6), Inches(2.5), Inches(7.5), Inches(0.07), RED)

add_text(s, Inches(0.6), Inches(2.0), Inches(8), Inches(0.5),
         'COSMAX  생산3팀', size=14, color=GRAY, bold=True)

add_text(s, Inches(0.6), Inches(2.7), Inches(12), Inches(1.6),
         '파우더 벌크 제조량 예측 프로그램', size=42, color=DARK, bold=True)
add_text(s, Inches(0.6), Inches(4.0), Inches(12), Inches(1),
         '운영 가이드 & ML 예측 FAQ', size=24, color=RED, bold=True)
add_text(s, Inches(0.6), Inches(6.5), Inches(8), Inches(0.5),
         '작성: 2026-04-27', size=13, color=GRAY)
add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.4),
         'github.com/djlee-cosmax/Powder-bulk-optimum-production-forecast',
         size=11, color=GRAY, font='Consolas')

# ============================================================
# Slide 2 — 목차
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header_bar(s, '목차', subtitle='Contents')

add_text(s, Inches(0.6), Inches(1.2), Inches(6), Inches(0.5),
         'PART 1.  프로그램 운영 가이드', size=20, color=BLUE, bold=True)
add_bullet_block(s, Inches(0.8), Inches(1.8), Inches(6), Inches(3), [
    '1.  프로그램 데이터 3종 구성',
    '2.  매주 월요일 데이터 업데이트 절차',
    '3.  명령어 동작 흐름 요약',
    '4.  자주 쓰는 명령어',
    '5.  트러블슈팅',
    '6.  깃 저장소 정보',
], size=15, accent=BLUE)

add_text(s, Inches(7.2), Inches(1.2), Inches(6), Inches(0.5),
         'PART 2.  ML 예측 FAQ', size=20, color=RED, bold=True)
add_bullet_block(s, Inches(7.4), Inches(1.8), Inches(6), Inches(3), [
    'Q1.  평균 로스율로 충분하지 않나?',
    'Q2.  ML은 어떻게 학습하나요?',
    'Q3.  오차란 무슨 의미인가요?',
], size=15, accent=RED)

# ============================================================
# PART 1 구분 슬라이드
# ============================================================
section_divider('PART 1', '프로그램 운영 가이드',
                '데이터 구성 · 주간 업데이트 워크플로 · 명령어 · 트러블슈팅', accent=BLUE)

# ============================================================
# Slide — 1. 프로그램 데이터 3종 구성
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header_bar(s, '1.  프로그램 데이터 3종 구성', accent=BLUE)

# 3개 카드
card_w = Inches(4.0)
card_h = Inches(4.5)
card_y = Inches(1.5)
gap = Inches(0.25)
total_w = card_w * 3 + gap * 2
card_x_start = (SW - total_w) // 2

cards = [
    ('①', '완제품 BOM 데이터', '매번 수동 업로드',
     '당일 작업할 제품의 BOM 정보\n→ 사용자가 매 작업 시 직접 업로드',
     ORANGE),
    ('②', '표준 대비 실적 데이터', '★ 주간 갱신 대상',
     'SAP 출력 xlsx\n→ 깃 저장소 자동 동기화\n팀원은 새로고침만 하면 최신',
     RED),
    ('③', '파우더 벌크 폐기 데이터', '★ 주간 갱신 대상',
     '화성·평택 시트 분리 xlsx\n→ 깃 저장소 자동 동기화\n팀원은 새로고침만 하면 최신',
     GREEN),
]

for i, (num, title, tag, desc, color) in enumerate(cards):
    x = card_x_start + (card_w + gap) * i
    add_rect(s, x, card_y, card_w, card_h, WHITE, line=LIGHT_GRAY)
    add_rect(s, x, card_y, card_w, Inches(0.6), color)
    add_text(s, x + Inches(0.3), card_y + Inches(0.05), Inches(1), Inches(0.5),
             num, size=22, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.0), card_y + Inches(0.05), Inches(3), Inches(0.5),
             title, size=15, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, x + Inches(0.3), card_y + Inches(0.85), card_w - Inches(0.6), Inches(0.4),
             tag, size=12, color=color, bold=True)
    add_text(s, x + Inches(0.3), card_y + Inches(1.4), card_w - Inches(0.6), Inches(2.8),
             desc, size=13, color=DARK)

# ============================================================
# Slide — 2. 매주 월요일 업데이트 절차
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header_bar(s, '2.  매주 월요일 데이터 업데이트 절차', accent=BLUE)

# 4단계 가로 흐름
step_y = Inches(1.3)
step_h = Inches(2.2)
step_w = Inches(2.95)
step_gap = Inches(0.2)
step_x_start = Inches(0.5)

steps = [
    ('STEP 1', '파일 받기',
     '바탕화면에 신규 파일 2개 수령\n(파일명에 날짜 자동 OK)', BLUE),
    ('STEP 2', 'WSL 터미널',
     'Windows Terminal → Ubuntu 또는\n시작 메뉴 "Ubuntu" 실행', BLUE),
    ('STEP 3', '명령 3줄 실행',
     '아래 명령 순서대로 실행\n(스크립트가 자동 처리)', RED),
    ('STEP 4', '완료',
     '팀원들 새로고침 시\n자동으로 최신 데이터 동기화', GREEN),
]

for i, (label, title, desc, color) in enumerate(steps):
    x = step_x_start + (step_w + step_gap) * i
    add_rect(s, x, step_y, step_w, step_h, WHITE, line=LIGHT_GRAY)
    add_rect(s, x, step_y, step_w, Inches(0.5), color)
    add_text(s, x, step_y + Inches(0.05), step_w, Inches(0.4),
             label, size=12, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), step_y + Inches(0.65), step_w - Inches(0.4), Inches(0.5),
             title, size=16, color=DARK, bold=True)
    add_text(s, x + Inches(0.2), step_y + Inches(1.15), step_w - Inches(0.4), Inches(1.0),
             desc, size=11, color=GRAY)

# 명령 코드 블록
add_text(s, Inches(0.5), Inches(3.7), Inches(12), Inches(0.4),
         'STEP 3 명령 (3줄)', size=13, color=DARK, bold=True)
code = """cd ~/cosmax/project2

bash update_data.sh \\
  "/mnt/c/Users/djlee/OneDrive - COSMAX/바탕 화면/표준 대비 실적 데이터_<날짜>.xlsx" \\
  "/mnt/c/Users/djlee/OneDrive - COSMAX/바탕 화면/파우더 벌크 폐기 데이터_<날짜>.xlsx"

git add data/ && git commit -m "데이터 업데이트 $(date +%Y-%m-%d)" && git push"""
add_code_block(s, Inches(0.5), Inches(4.15), Inches(12.3), Inches(2.9), code, size=11)

# ============================================================
# Slide — 3. 명령어 동작 흐름 요약
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header_bar(s, '3.  명령어 동작 흐름 요약', accent=BLUE)

# 좌: update_data.sh
add_rect(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(2.7), LIGHT_GRAY)
add_rect(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.5), BLUE)
add_text(s, Inches(0.7), Inches(1.35), Inches(6), Inches(0.4),
         '🔧  update_data.sh 가 하는 일', size=14, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
add_bullet_block(s, Inches(0.7), Inches(2.0), Inches(5.8), Inches(2.0), [
    '신규 xlsx 2개를 ~/cosmax/project2/data/ 에 복사',
    '파일명을 standard_perf.xlsx / disposal.xlsx 로 통일',
    'data/manifest.json 의 updated 날짜를 오늘로 갱신',
], size=13, accent=BLUE)

# 우: 자동 동기화
add_rect(s, Inches(6.85), Inches(1.3), Inches(6.0), Inches(2.7), LIGHT_GRAY)
add_rect(s, Inches(6.85), Inches(1.3), Inches(6.0), Inches(0.5), GREEN)
add_text(s, Inches(7.05), Inches(1.35), Inches(6), Inches(0.4),
         '🔄  팀원 측 자동 동기화 동작', size=14, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
add_bullet_block(s, Inches(7.05), Inches(2.0), Inches(5.8), Inches(2.0), [
    '페이지 로드 시 manifest.json 의 updated 비교',
    '본인 IndexedDB 캐시의 serverVersion 과 대조',
    '다르면 신규 다운로드 → 파싱 → 캐시 갱신',
    '같으면 캐시 즉시 사용 (빠름)',
], size=13, accent=GREEN)

# 하: 흐름도
add_text(s, Inches(0.5), Inches(4.3), Inches(12), Inches(0.4),
         '전체 흐름', size=13, color=DARK, bold=True)
flow_y = Inches(4.85)
flow_h = Inches(1.6)
flow_steps = [
    ('xlsx 받기', '바탕화면', BLUE),
    ('update_data.sh', '복사 + manifest 갱신', RED),
    ('git push', '깃에 배포', RED),
    ('팀원 새로고침', '자동 동기화', GREEN),
    ('최신 데이터 사용', '완료', GREEN),
]
flow_w = Inches(2.3)
flow_gap = Inches(0.15)
flow_x = Inches(0.6)
for i, (title, sub, color) in enumerate(flow_steps):
    add_rect(s, flow_x, flow_y, flow_w, flow_h, color)
    add_text(s, flow_x, flow_y + Inches(0.3), flow_w, Inches(0.5),
             title, size=14, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, flow_x, flow_y + Inches(0.85), flow_w, Inches(0.5),
             sub, size=11, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    flow_x += flow_w + flow_gap
    if i < len(flow_steps) - 1:
        # 화살표
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   flow_x - flow_gap - Pt(2),
                                   flow_y + flow_h // 2 - Inches(0.1),
                                   flow_gap, Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK
        arrow.line.fill.background()

# ============================================================
# Slide — 4. 자주 쓰는 명령어
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header_bar(s, '4.  자주 쓰는 명령어', accent=BLUE)

cmd_groups = [
    ('🖥️  서버 시작 (PC 부팅 후 1회)',
     'cd ~/cosmax && node server.js', BLUE),
    ('📤  수동 업로드 (특정 일자 데이터 일시 점검)',
     '브라우저 화면 → "표준 대비 실적 데이터" / "벌크 폐기 데이터" 업로드 버튼 클릭\n※ 다음 새로고침 시 서버 데이터로 덮어쓰여짐 (정상 동작)', ORANGE),
    ('🔄  캐시 초기화 (전체 다시 받기)',
     '화면의 "전체 초기화" 버튼 → 새로고침', GREEN),
]

y = Inches(1.4)
for title, cmd, color in cmd_groups:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.5), color)
    add_text(s, Inches(0.7), y + Inches(0.05), Inches(12), Inches(0.4),
             title, size=14, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if '\n' in cmd:
        add_text(s, Inches(0.7), y + Inches(0.7), Inches(12), Inches(1.0),
                 cmd, size=12, color=DARK)
        y += Inches(2.0)
    else:
        add_code_block(s, Inches(0.5), y + Inches(0.55), Inches(12.3), Inches(0.6), cmd, size=12)
        y += Inches(1.5)

# ============================================================
# Slide — 5. 트러블슈팅
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header_bar(s, '5.  트러블슈팅', accent=BLUE)

trouble_items = [
    ('git push 시 인증 요구', 'GitHub 토큰/계정 입력 (1회만 저장하면 다음부터 자동)'),
    ('데이터가 갱신 안 됨', '브라우저에서 Ctrl + Shift + R (강력 새로고침)\n그래도 안 되면 "전체 초기화" 클릭'),
    ('서버 접속 안 됨 (localhost:8001)', 'WSL 터미널에서  cd ~/cosmax && node server.js  다시 실행'),
    ('폐기 데이터 시트 이름', '반드시 "화성", "평택" 단어가 시트명에 포함되어야 매칭됨\n예: "화성25년", "평택4월" 등은 OK'),
]

y = Inches(1.4)
for title, desc in trouble_items:
    add_rect(s, Inches(0.5), y, Inches(0.15), Inches(1.0), RED)
    add_text(s, Inches(0.85), y + Inches(0.05), Inches(11.5), Inches(0.4),
             '▶  ' + title, size=15, color=DARK, bold=True)
    add_text(s, Inches(1.05), y + Inches(0.5), Inches(11.5), Inches(0.55),
             desc, size=12, color=GRAY)
    y += Inches(1.18)

# ============================================================
# Slide — 6. 깃 저장소 정보
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header_bar(s, '6.  깃 저장소 정보', accent=BLUE)

info_items = [
    ('🌐 저장소 URL',
     'https://github.com/djlee-cosmax/Powder-bulk-optimum-production-forecast'),
    ('📁 로컬 위치',
     '~/cosmax/project2  (WSL Ubuntu 내부)'),
    ('💾 데이터 파일 위치',
     '~/cosmax/project2/data/  (standard_perf.xlsx, disposal.xlsx, manifest.json)'),
    ('🔧 헬퍼 스크립트',
     '~/cosmax/project2/update_data.sh'),
]

y = Inches(1.6)
for label, val in info_items:
    add_rect(s, Inches(0.7), y, Inches(11.9), Inches(1.0), LIGHT_GRAY)
    add_text(s, Inches(1.0), y + Inches(0.1), Inches(4), Inches(0.4),
             label, size=14, color=DARK, bold=True)
    add_text(s, Inches(1.0), y + Inches(0.55), Inches(11), Inches(0.4),
             val, size=12, color=BLUE, font='Consolas')
    y += Inches(1.15)

# ============================================================
# PART 2 구분 슬라이드
# ============================================================
section_divider('PART 2', 'ML 예측 FAQ',
                '왜 ML을 쓰는지 · 어떻게 학습하는지 · 오차의 의미', accent=RED)

# ============================================================
# Slide — Q1: 평균 로스율 vs ML
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header_bar(s, 'Q1.  평균 로스율이면 되는데 ML은 왜 쓰나요?', accent=RED)

add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
         'A.  평균은 모든 상황을 하나의 숫자로 퉁치는 것입니다.',
         size=18, color=DARK, bold=True)

# 비교 박스 - 같은 벌크 다른 수량
add_text(s, Inches(0.5), Inches(2.0), Inches(12), Inches(0.4),
         '예시.  같은 벌크라도 수량에 따라 로스율 차이가 큽니다.',
         size=14, color=GRAY, bold=True)

# 박스 3개 비교
box_y = Inches(2.55)
box_h = Inches(1.6)
box_w = Inches(4.0)

# 박스1 - 소량
add_rect(s, Inches(0.5), box_y, box_w, box_h, RGBColor(0xff, 0xf4, 0xf4),
         line=RED)
add_text(s, Inches(0.7), box_y + Inches(0.15), Inches(3.6), Inches(0.4),
         '🔻 소량 생산 (5,000g)', size=14, color=RED, bold=True)
add_text(s, Inches(0.7), box_y + Inches(0.65), Inches(3.6), Inches(0.4),
         '로스율 ≈ 55%', size=22, color=RED, bold=True)
add_text(s, Inches(0.7), box_y + Inches(1.15), Inches(3.6), Inches(0.4),
         '평균 27% 적용 시 → 벌크 부족',
         size=11, color=GRAY)

# 박스2 - 평균 (중앙)
add_rect(s, Inches(4.65), box_y, box_w, box_h, LIGHT_GRAY,
         line=GRAY)
add_text(s, Inches(4.85), box_y + Inches(0.15), Inches(3.6), Inches(0.4),
         '➖ 단순 평균', size=14, color=GRAY, bold=True)
add_text(s, Inches(4.85), box_y + Inches(0.65), Inches(3.6), Inches(0.4),
         '로스율 = 27%', size=22, color=GRAY, bold=True)
add_text(s, Inches(4.85), box_y + Inches(1.15), Inches(3.6), Inches(0.4),
         '소량은 부족 / 대량은 잉여 발생',
         size=11, color=GRAY)

# 박스3 - 대량
add_rect(s, Inches(8.8), box_y, box_w, box_h, RGBColor(0xf4, 0xfa, 0xf4),
         line=GREEN)
add_text(s, Inches(9.0), box_y + Inches(0.15), Inches(3.6), Inches(0.4),
         '🔺 대량 생산 (50,000g)', size=14, color=GREEN, bold=True)
add_text(s, Inches(9.0), box_y + Inches(0.65), Inches(3.6), Inches(0.4),
         '로스율 ≈ 8%', size=22, color=GREEN, bold=True)
add_text(s, Inches(9.0), box_y + Inches(1.15), Inches(3.6), Inches(0.4),
         '평균 27% 적용 시 → 벌크 잉여',
         size=11, color=GRAY)

# ML 강조 박스
ml_box_y = Inches(4.45)
add_rect(s, Inches(0.5), ml_box_y, Inches(12.3), Inches(2.5), RED)
add_text(s, Inches(0.8), ml_box_y + Inches(0.2), Inches(11), Inches(0.5),
         '✨  ML이 하는 일', size=16, color=WHITE, bold=True)
add_text(s, Inches(0.8), ml_box_y + Inches(0.85), Inches(11.5), Inches(0.6),
         '수량 · 설비 · 제품 유형 · 공장 · 시기 등 여러 조건을 동시에 고려하여',
         size=14, color=WHITE)
add_text(s, Inches(0.8), ml_box_y + Inches(1.25), Inches(11.5), Inches(0.6),
         '그 상황에 맞는 로스율을 예측합니다.',
         size=14, color=WHITE)
add_text(s, Inches(0.8), ml_box_y + Inches(1.85), Inches(11.5), Inches(0.5),
         '실제로 오차가  20% → 11% 로 약 45% 개선되었습니다.',
         size=15, color=WHITE, bold=True)

# ============================================================
# Slide — Q2: 어떻게 학습?
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header_bar(s, 'Q2.  어떤 식으로 학습하나요?', accent=RED)

add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
         'A.  과거 실적 데이터 3만 건에서 패턴을 찾는 것입니다.',
         size=18, color=DARK, bold=True)

# 사람 비유
add_rect(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.0), LIGHT_GRAY)
add_text(s, Inches(0.7), Inches(2.1), Inches(12), Inches(0.4),
         '👨‍💼  사람으로 비유하면', size=14, color=DARK, bold=True)
add_text(s, Inches(0.7), Inches(2.55), Inches(12), Inches(0.5),
         '신입사원:  "로스율은 대충 25~30% 잡으세요"',
         size=14, color=GRAY)
add_text(s, Inches(0.7), Inches(3.05), Inches(12), Inches(0.5),
         '경험 쌓인 후:  "이 벌크는 소량이면 많이 잡아야 하고,',
         size=14, color=DARK)
add_text(s, Inches(0.7), Inches(3.5), Inches(12), Inches(0.5),
         '              이 설비에서는 좀 덜 나와요"',
         size=14, color=DARK)

# ML 흐름
add_text(s, Inches(0.5), Inches(4.3), Inches(12), Inches(0.5),
         '🤖  ML이 그 "경험 축적"을 데이터로 합니다.',
         size=15, color=RED, bold=True)

# 3단계
ml_steps = [
    ('STEP 1', '데이터 학습', '3만 건 실적에서\n"이 조건일 때 로스율이 얼마였는지"\n패턴을 학습'),
    ('STEP 2', '신규 주문 입력', '새 주문이 들어오면\n수량·설비·유형·공장·시기 등\n조건을 받음'),
    ('STEP 3', '예측', '가장 비슷한 과거 사례들을\n종합해서 그 상황에 맞는\n로스율 예측'),
]
y = Inches(5.0)
w = Inches(4.0)
h = Inches(2.0)
gap = Inches(0.15)
x_start = (SW - w * 3 - gap * 2) // 2
for i, (lbl, ttl, dsc) in enumerate(ml_steps):
    x = x_start + (w + gap) * i
    add_rect(s, x, y, w, h, WHITE, line=RED)
    add_rect(s, x, y, w, Inches(0.4), RED)
    add_text(s, x, y + Inches(0.05), w, Inches(0.3),
             lbl, size=11, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), y + Inches(0.5), w - Inches(0.4), Inches(0.4),
             ttl, size=14, color=DARK, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.95), w - Inches(0.4), Inches(1.0),
             dsc, size=11, color=GRAY)

# ============================================================
# Slide — Q3: 오차의 의미
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header_bar(s, 'Q3.  오차(MAE)는 무슨 의미인가요?', accent=RED)

add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
         'A.  예측한 로스율과 실제 로스율의 차이입니다.',
         size=18, color=DARK, bold=True)

# 단일 사례 비교
add_text(s, Inches(0.5), Inches(2.0), Inches(12), Inches(0.4),
         '예시.  실제 로스율이 20% 인 건이 있을 때',
         size=14, color=GRAY, bold=True)

# 박스 2개
add_rect(s, Inches(0.7), Inches(2.5), Inches(5.85), Inches(1.3), LIGHT_GRAY,
         line=GRAY)
add_text(s, Inches(0.9), Inches(2.6), Inches(5.5), Inches(0.4),
         '➖  평균 방식 예측', size=13, color=GRAY, bold=True)
add_text(s, Inches(0.9), Inches(3.0), Inches(5.5), Inches(0.4),
         '예측: 27%   →   실제: 20%', size=14, color=DARK)
add_text(s, Inches(0.9), Inches(3.4), Inches(5.5), Inches(0.4),
         '오차 = 7%p', size=18, color=GRAY, bold=True)

add_rect(s, Inches(6.75), Inches(2.5), Inches(5.85), Inches(1.3),
         RGBColor(0xff, 0xf4, 0xf4), line=RED)
add_text(s, Inches(6.95), Inches(2.6), Inches(5.5), Inches(0.4),
         '✨  ML 예측', size=13, color=RED, bold=True)
add_text(s, Inches(6.95), Inches(3.0), Inches(5.5), Inches(0.4),
         '예측: 22%   →   실제: 20%', size=14, color=DARK)
add_text(s, Inches(6.95), Inches(3.4), Inches(5.5), Inches(0.4),
         '오차 = 2%p', size=18, color=RED, bold=True)

# 전체 효과
add_text(s, Inches(0.5), Inches(4.1), Inches(12), Inches(0.4),
         '3만 건 전체 평균',
         size=14, color=GRAY, bold=True)

# 거대한 비교 막대
add_rect(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(0.9),
         LIGHT_GRAY)
add_text(s, Inches(0.9), Inches(4.7), Inches(4), Inches(0.4),
         '평균 사용', size=13, color=GRAY, bold=True)
# 막대 (비례 - 20.1)
bar_max_w = Inches(7.0)
add_rect(s, Inches(3.5), Inches(4.85), bar_max_w, Inches(0.4), GRAY)
add_text(s, Inches(11.0), Inches(4.85), Inches(1.6), Inches(0.4),
         '오차 20.1%', size=14, color=WHITE, bold=True,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.9), RGBColor(0xff, 0xf4, 0xf4))
add_text(s, Inches(0.9), Inches(5.7), Inches(4), Inches(0.4),
         'ML 사용', size=13, color=RED, bold=True)
# 막대 비례 (11.14 / 20.1 ≈ 0.554)
ml_bar_w = int(bar_max_w * 11.14 / 20.1)
add_rect(s, Inches(3.5), Inches(5.85), ml_bar_w, Inches(0.4), RED)
add_text(s, Inches(3.5) + ml_bar_w + Inches(0.1), Inches(5.85), Inches(2), Inches(0.4),
         '오차 11.14%', size=14, color=RED, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# 효과 강조
add_rect(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.7), GREEN)
add_text(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.6),
         '✓  오차 ↓  =  제조량 정확도 ↑  =  잔량 / 폐기 감소',
         size=15, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# 마지막 슬라이드 — 정리
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
add_rect(s, 0, Inches(0.5), Inches(0.3), SH - Inches(1), RED)

add_text(s, Inches(0.7), Inches(0.7), Inches(12), Inches(0.5),
         'COSMAX  생산3팀', size=13, color=GRAY, bold=True)
add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(1),
         '한 줄 요약', size=36, color=DARK, bold=True)

add_rect(s, Inches(0.7), Inches(2.5), Inches(12), Inches(1.2), LIGHT_GRAY)
add_text(s, Inches(0.9), Inches(2.6), Inches(12), Inches(0.5),
         'Part 1 운영', size=14, color=BLUE, bold=True)
add_text(s, Inches(0.9), Inches(3.05), Inches(12), Inches(0.6),
         '월요일에 명령어 3줄로 데이터 갱신 → 팀원은 새로고침 한 번이면 끝',
         size=15, color=DARK)

add_rect(s, Inches(0.7), Inches(3.9), Inches(12), Inches(1.2), RGBColor(0xff, 0xf4, 0xf4))
add_text(s, Inches(0.9), Inches(4.0), Inches(12), Inches(0.5),
         'Part 2 ML', size=14, color=RED, bold=True)
add_text(s, Inches(0.9), Inches(4.45), Inches(12), Inches(0.6),
         'ML은 상황별 로스율을 예측해 평균 대비 오차를 45% 줄였습니다.',
         size=15, color=DARK)

add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
         '문의 / 이슈',
         size=12, color=GRAY, bold=True)
add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         'github.com/djlee-cosmax/Powder-bulk-optimum-production-forecast',
         size=11, color=BLUE, font='Consolas')

# ===== 저장 =====
out_path = "/mnt/c/Users/djlee/OneDrive - COSMAX/바탕 화면/파우더 벌크 제조량 예측 프로그램 종합 안내.pptx"
prs.save(out_path)
print(f"PPT 생성 완료: {out_path}")
print(f"슬라이드 수: {len(prs.slides)}장")
