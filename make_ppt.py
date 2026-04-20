#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

RED = RGBColor(0xC8, 0x10, 0x2E)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MID = RGBColor(0x44, 0x44, 0x44)
LIGHT = RGBColor(0x88, 0x88, 0x88)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
BG_INFO = RGBColor(0xF0, 0xF4, 0xFF)
BLUE = RGBColor(0x4E, 0x79, 0xA7)
FONT = "Malgun Gothic"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]


def set_run(run, text, size=14, bold=False, color=DARK, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return tb


def add_line(slide, x, y, w, h, color=RED):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_box(slide, x, y, w, h, fill=BG_LIGHT, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.shadow.inherit = False
    return sh


def slide_header(slide, num, title):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(0.6), Inches(0.6), f"{num:02d}",
                size=32, bold=True, color=RED)
    add_textbox(slide, Inches(1.25), Inches(0.38), Inches(11), Inches(0.6), title,
                size=26, bold=True, color=DARK)
    add_line(slide, Inches(0.5), Inches(1.1), Inches(12.33), Emu(38100), color=RED)


def slide_footer(slide, pagenum, total):
    add_textbox(slide, Inches(0.5), Inches(7.08), Inches(7), Inches(0.3),
                "COSMAX 생산3팀 | 파우더 벌크 제조량 예측 프로그램",
                size=9, color=LIGHT)
    add_textbox(slide, Inches(11.8), Inches(7.08), Inches(1), Inches(0.3),
                f"{pagenum} / {total}", size=9, color=LIGHT, align=PP_ALIGN.RIGHT)


def add_bullet(slide, x, y, w, h, items, size=14, line_spacing=1.35):
    """items: list of (label_bold, description) or plain string"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if isinstance(it, tuple):
            label, desc = it
            r1 = p.add_run(); set_run(r1, "• ", size=size, bold=True, color=RED)
            r2 = p.add_run(); set_run(r2, label, size=size, bold=True, color=DARK)
            if desc:
                r3 = p.add_run(); set_run(r3, " — " + desc, size=size, bold=False, color=MID)
        else:
            r1 = p.add_run(); set_run(r1, "• ", size=size, bold=True, color=RED)
            r2 = p.add_run(); set_run(r2, it, size=size, bold=False, color=DARK)
    return tb


def add_table(slide, x, y, w, h, data, header_fill=RED, col_widths=None):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.06)
            cell.margin_bottom = Inches(0.06)
            text = str(data[r][c])
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                set_run(run, text, size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r % 2 == 1 else RGBColor(0xFA, 0xFA, 0xFA)
                set_run(run, text, size=11, bold=False, color=DARK)
    return tbl


TOTAL = 10

# ===================== SLIDE 1: 표지 =====================
s = prs.slides.add_slide(blank)
# 상단 빨간 바
add_line(s, 0, 0, SW, Inches(0.12), color=RED)
# 중앙 타이틀
add_textbox(s, Inches(1), Inches(2.4), Inches(11.3), Inches(0.6), "COSMAX 생산3팀",
            size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(1), Inches(3.1), Inches(11.3), Inches(1.2),
            "파우더 벌크 제조량 예측 프로그램",
            size=40, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_line(s, Inches(6.17), Inches(4.35), Inches(1), Inches(0.04), color=RED)
add_textbox(s, Inches(1), Inches(4.55), Inches(11.3), Inches(0.5),
            "SAP 실적 데이터 + 머신러닝 기반 최적 제조량 예측 시스템",
            size=18, bold=False, color=MID, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(1), Inches(6.4), Inches(11.3), Inches(0.4),
            "버전 2.0  |  2026년 4월",
            size=12, color=LIGHT, align=PP_ALIGN.CENTER)
add_line(s, 0, Inches(7.38), SW, Inches(0.12), color=RED)

# ===================== SLIDE 2: 프로그램 개요 =====================
s = prs.slides.add_slide(blank)
slide_header(s, 1, "프로그램 개요")

add_textbox(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
            "파우더 제품의 '벌크 최적 제조량'을 예측하는 도구",
            size=18, bold=True, color=DARK)
add_textbox(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.9),
            "완제품 BOM 데이터와 SAP 실적 데이터를 기반으로, 각 벌크의 이론 필요량에 과거 로스율을 반영하여\n실제로 제조해야 할 최적 수량을 산출합니다.",
            size=14, color=MID)

# 입출력 다이어그램
box_y = Inches(3.3)
box_h = Inches(1.1)
# 입력 박스
add_box(s, Inches(0.7), box_y, Inches(3.5), box_h, fill=BG_INFO, line=BLUE)
add_textbox(s, Inches(0.7), box_y, Inches(3.5), Inches(0.35),
            "INPUT  |  입력 데이터", size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.7), box_y+Inches(0.35), Inches(3.5), Inches(0.75),
            "• 완제품 BOM\n• SAP 표준 대비 실적\n• 벌크 폐기 데이터", size=12, color=DARK, align=PP_ALIGN.CENTER)

# 화살표
add_textbox(s, Inches(4.3), box_y+Inches(0.3), Inches(0.7), Inches(0.5),
            "▶", size=30, bold=True, color=RED, align=PP_ALIGN.CENTER)

# 처리 박스
add_box(s, Inches(5.0), box_y, Inches(3.5), box_h, fill=BG_LIGHT, line=RED)
add_textbox(s, Inches(5.0), box_y, Inches(3.5), Inches(0.35),
            "PROCESS  |  예측 계산", size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(5.0), box_y+Inches(0.35), Inches(3.5), Inches(0.75),
            "• 최근 5건 가중평균 로스율\n• 신뢰도 점수 산출\n• ML 예측 비교", size=12, color=DARK, align=PP_ALIGN.CENTER)

# 화살표
add_textbox(s, Inches(8.6), box_y+Inches(0.3), Inches(0.7), Inches(0.5),
            "▶", size=30, bold=True, color=RED, align=PP_ALIGN.CENTER)

# 출력 박스
add_box(s, Inches(9.3), box_y, Inches(3.5), box_h, fill=BG_INFO, line=BLUE)
add_textbox(s, Inches(9.3), box_y, Inches(3.5), Inches(0.35),
            "OUTPUT  |  결과", size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(9.3), box_y+Inches(0.35), Inches(3.5), Inches(0.75),
            "• 벌크별 최적 제조량\n• 신뢰도 + 상세 이력\n• 엑셀/CSV 다운로드", size=12, color=DARK, align=PP_ALIGN.CENTER)

# 핵심 한 줄
add_box(s, Inches(0.7), Inches(5.2), Inches(12.1), Inches(1.3), fill=RGBColor(0xFF, 0xF5, 0xF5))
add_textbox(s, Inches(0.9), Inches(5.35), Inches(11.7), Inches(0.4),
            "핵심 가치", size=12, bold=True, color=RED)
add_textbox(s, Inches(0.9), Inches(5.75), Inches(11.7), Inches(0.7),
            "SAP 표준 대비 투입률 데이터를 활용해 '이론 필요량'에 로스율을 자동 반영 → 현장에서 실측 없이도 적정 제조량 도출",
            size=13, color=DARK)

slide_footer(s, 2, TOTAL)

# ===================== SLIDE 3: 주요 기능 =====================
s = prs.slides.add_slide(blank)
slide_header(s, 2, "주요 기능")

features = [
    ("BOM 구조 분석", "완제품 → 성형물 → 벌크 구조를 트리 형태로 시각화"),
    ("벌크별 필요 제조량 계산", "발주 수량 입력 시 벌크별 최적 제조량 자동 산출"),
    ("로스율 기반 예측", "SAP 실적 데이터 최근 5건 가중평균으로 로스율 산출 (이상치 영향 감소)"),
    ("신뢰도 점수 표시", "표본수·편차·최신성 기반 0~100점 평가 및 등급 배지"),
    ("ML(머신러닝) 예측 비교", "42,408건 학습 모델로 제조량 구간별 로스율 교차 검증"),
    ("폐기 데이터 보정", "폐기된 벌크량을 차감하여 실제 사용량 기준으로 계산"),
    ("데이터 자동 저장", "IndexedDB에 SAP·폐기 데이터 저장, 페이지 재접속 시 자동 복원"),
    ("결과 내보내기", "엑셀(XLSX) / CSV 파일 다운로드 지원 (신뢰도 컬럼 포함 14열)"),
]
add_bullet(s, Inches(0.7), Inches(1.4), Inches(12), Inches(5.5), features, size=15, line_spacing=1.45)

slide_footer(s, 3, TOTAL)

# ===================== SLIDE 4: 필요 데이터 3종 =====================
s = prs.slides.add_slide(blank)
slide_header(s, 3, "필요 데이터 3종")

data = [
    ["데이터", "버튼 색상", "필수 여부", "설명"],
    ["완제품 BOM 데이터", "보라색", "필수", "SAP에서 추출한 BOM (TXT / CSV)"],
    ["표준 대비 실적 데이터", "파란색", "필수", "SAP 표준소요량 대비 실제 투입량 이력 (CSV / XLSX)"],
    ["벌크 폐기 데이터", "초록색", "선택", "폐기 처리된 벌크 수량 (CSV / XLSX, 화성/평택 시트)"],
]
add_table(s, Inches(0.7), Inches(1.5), Inches(12), Inches(2.3), data,
          col_widths=[Inches(3), Inches(1.8), Inches(1.5), Inches(5.7)])

# 참고 박스
add_box(s, Inches(0.7), Inches(4.3), Inches(12), Inches(1.2), fill=BG_INFO, line=BLUE)
add_textbox(s, Inches(0.9), Inches(4.4), Inches(11.6), Inches(0.3),
            "참고", size=12, bold=True, color=BLUE)
add_textbox(s, Inches(0.9), Inches(4.7), Inches(11.6), Inches(0.8),
            "• 폐기 데이터가 없어도 예측은 가능 — 다만 함께 업로드하면 실제 사용량을 보정하여 더 정확한 로스율 산출\n• 업로드 순서는 자유 / '벌크 필요량 계산' 버튼 누르기 전에 표준 대비 실적 데이터만 있으면 됨",
            size=12, color=DARK)

# 현재 데이터 규모
add_textbox(s, Inches(0.7), Inches(5.7), Inches(12.3), Inches(0.4),
            "2026년 4월 기준 데이터 규모", size=14, bold=True, color=RED)
add_bullet(s, Inches(0.7), Inches(6.1), Inches(12.3), Inches(1.0), [
    ("표준 대비 실적", "46,354건 (2023.01 ~ 2026.04, 3년 3개월)"),
    ("벌크 폐기", "366건 (화성 134 / 평택 232, 최근 1개월)"),
], size=13, line_spacing=1.35)

slide_footer(s, 4, TOTAL)

# ===================== SLIDE 5: 예측 계산 공식 =====================
s = prs.slides.add_slide(blank)
slide_header(s, 4, "예측 계산 방식")

# Step 1
add_textbox(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.4),
            "Step 1. 이론 필요량 산출", size=15, bold=True, color=RED)
add_box(s, Inches(0.9), Inches(1.85), Inches(11.6), Inches(0.7), fill=BG_LIGHT)
add_textbox(s, Inches(0.9), Inches(1.95), Inches(11.6), Inches(0.5),
            "이론 필요량(g) = 성형물 필요 수량(ea) × 벌크 투입량(g)",
            size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# Step 2
add_textbox(s, Inches(0.7), Inches(2.85), Inches(12), Inches(0.4),
            "Step 2. 로스율 산출 (최근 5건 가중평균)", size=15, bold=True, color=RED)
add_box(s, Inches(0.9), Inches(3.35), Inches(11.6), Inches(0.7), fill=BG_LIGHT)
add_textbox(s, Inches(0.9), Inches(3.45), Inches(11.6), Inches(0.5),
            "로스율(%) = (실제 투입량 − 표준소요량) ÷ 표준소요량 × 100",
            size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_bullet(s, Inches(0.9), Inches(4.15), Inches(11.6), Inches(1.2), [
    "과거 이력 최신순 정렬 → 유효한 최근 5건을 '가중평균' (최신일수록 큰 가중치: 5,4,3,2,1)",
    "표준소요량 3,000g 이하는 제외 (소량 제조 시 로스율 왜곡 방지)",
    "폐기 데이터가 있으면 실제 투입량에서 폐기량 차감 후 로스율 계산",
], size=12, line_spacing=1.3)

# Step 3
add_textbox(s, Inches(0.7), Inches(5.55), Inches(12), Inches(0.4),
            "Step 3. 최적 제조량 산출", size=15, bold=True, color=RED)
add_box(s, Inches(0.9), Inches(6.05), Inches(11.6), Inches(0.7), fill=BG_LIGHT)
add_textbox(s, Inches(0.9), Inches(6.15), Inches(11.6), Inches(0.5),
            "최적 제조량(g) = 이론 필요량 × (1 + 로스율 ÷ 100)",
            size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.9), Inches(6.8), Inches(11.6), Inches(0.3),
            "예시: 57,000g × (1 + 10.59/100) = 63,036g",
            size=11, color=MID, align=PP_ALIGN.CENTER)

slide_footer(s, 5, TOTAL)

# ===================== SLIDE 6: 신뢰도 점수 =====================
s = prs.slides.add_slide(blank)
slide_header(s, 5, "신뢰도 점수 시스템 (0~100점)")

add_textbox(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.5),
            "각 예측에 대해 얼마나 믿을 수 있는지를 직관적으로 확인 — 표본수·편차·최신성 3요소로 평가",
            size=13, color=MID)

data = [
    ["항목 (만점)", "기준", "점수"],
    ["표본 수 (40점)", "5건 이상", "40점"],
    ["", "3~4건", "25점"],
    ["", "1~2건", "10점"],
    ["편차 (40점)", "표준편차 3% 미만", "40점"],
    ["", "표준편차 7% 미만", "25점"],
    ["", "표준편차 15% 미만", "10점"],
    ["최신성 (20점)", "최신 이력 30일 이내", "20점"],
    ["", "최신 이력 90일 이내", "10점"],
    ["", "최신 이력 180일 이내", "5점"],
]
add_table(s, Inches(0.7), Inches(2.0), Inches(6.5), Inches(4.6), data,
          col_widths=[Inches(2), Inches(3), Inches(1.5)])

# 등급 표시
add_textbox(s, Inches(7.5), Inches(2.0), Inches(5.3), Inches(0.4),
            "등급 및 취급 기준", size=13, bold=True, color=RED)

grade_data = [
    ["등급", "점수", "취급"],
    ["높음", "80+", "그대로 활용 가능"],
    ["보통", "50~79", "활용 가능하나 주의"],
    ["낮음", "30~49", "참고용 · 실측 권장"],
    ["매우 낮음", "~29", "실측 필수"],
]
add_table(s, Inches(7.5), Inches(2.45), Inches(5.3), Inches(2.5), grade_data,
          col_widths=[Inches(1.4), Inches(1.2), Inches(2.7)])

# 시각 강조
add_box(s, Inches(7.5), Inches(5.15), Inches(5.3), Inches(1.45), fill=RGBColor(0xFF, 0xF5, 0xF5))
add_textbox(s, Inches(7.7), Inches(5.25), Inches(5.0), Inches(0.3),
            "시각적 강조", size=12, bold=True, color=RED)
add_textbox(s, Inches(7.7), Inches(5.55), Inches(5.0), Inches(1.0),
            "• 낮음 → 옅은 주황 배경\n• 매우 낮음 → 옅은 빨강 배경 (실측 필수)\n• 툴팁에 표본수·편차·최신성 상세 표시",
            size=11, color=DARK)

slide_footer(s, 6, TOTAL)

# ===================== SLIDE 7: ML 예측 =====================
s = prs.slides.add_slide(blank)
slide_header(s, 6, "ML(머신러닝) 예측 방식")

# 학습 개요
add_textbox(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.5),
            "SAP 실적 데이터를 Gradient Boosting 알고리즘으로 학습 — 벌크코드별 + 제조량 구간별 예측",
            size=13, color=MID)

# 지표 박스 (3개)
kpi_y = Inches(2.0)
kpi_h = Inches(1.3)
for i, (val, label, color) in enumerate([
    ("42,408건", "학습 데이터", RED),
    ("6,131종", "예측 벌크 수", BLUE),
    ("11.91%", "테스트 MAE", RED),
]):
    x = Inches(0.7 + i * 4.1)
    add_box(s, x, kpi_y, Inches(3.9), kpi_h, fill=BG_LIGHT)
    add_textbox(s, x, kpi_y+Inches(0.15), Inches(3.9), Inches(0.6), val,
                size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(s, x, kpi_y+Inches(0.85), Inches(3.9), Inches(0.4), label,
                size=12, color=MID, align=PP_ALIGN.CENTER)

# 학습 변수
add_textbox(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.4),
            "학습에 사용되는 7개 변수", size=14, bold=True, color=RED)
var_data = [
    ["변수", "설명"],
    ["실적수량", "실제 생산된 성형물 수량 (제조 규모)"],
    ["관리유형", "제품의 관리 유형 분류"],
    ["작업장", "공장·설비가 결합된 작업장 정보 (예: 화성 자동 파우더 성형기2)"],
    ["벌크 평균 로스율", "해당 벌크의 과거 평균 로스율"],
    ["벌크 로스율 표준편차", "로스율 변동성"],
    ["벌크 로스율 중간값", "이상치에 덜 민감한 대표값"],
    ["평균 단위당 표준소요량", "성형물 1개당 벌크 투입량 평균"],
]
add_table(s, Inches(0.7), Inches(4.15), Inches(12), Inches(2.75), var_data,
          col_widths=[Inches(3.5), Inches(8.5)])

slide_footer(s, 7, TOTAL)

# ===================== SLIDE 8: SAP vs ML =====================
s = prs.slides.add_slide(blank)
slide_header(s, 7, "SAP 방식 vs ML 방식 비교")

data = [
    ["구분", "SAP 방식 (기본)", "ML 방식 (참고)"],
    ["로스율 기준", "최근 5건의 가중평균 로스율", "전체 이력 + 다양한 조건 학습 결과"],
    ["제조량 반영", "반영하지 않음", "제조량 구간별 다른 로스율 적용"],
    ["이력 1건일 때", "그 1건의 로스율 (신뢰도 낮음)", "작업장·관리유형·제조량 등 조건 종합"],
    ["이력 없을 때", "- 표시 (예측 불가)", "- 표시 (학습 데이터 없음)"],
    ["장점", "최근 현장 상황 즉시 반영\n이상치 1건 영향 감소", "제조량·작업장·관리유형 조건 반영"],
    ["활용", "기본 예측값", "교차 검증용 참고값"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.3), data,
          col_widths=[Inches(2.3), Inches(5), Inches(5)])

# 결론 박스
add_box(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0), fill=RGBColor(0xFF, 0xF5, 0xF5))
add_textbox(s, Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.3),
            "결론", size=12, bold=True, color=RED)
add_textbox(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.6),
            "SAP 방식을 기본값으로 사용하고, ML 예측은 교차 검증용 참고값으로 활용. 두 예측의 차이가 크면(5%p↑ 노랑 / 10%p↑ 빨강) 실측 검토 필요.",
            size=12, color=DARK)

slide_footer(s, 8, TOTAL)

# ===================== SLIDE 9: 결과 화면 기능 =====================
s = prs.slides.add_slide(blank)
slide_header(s, 8, "결과 화면 기능")

add_textbox(s, Inches(0.7), Inches(1.4), Inches(6), Inches(0.4),
            "결과 테이블 항목", size=14, bold=True, color=RED)
add_bullet(s, Inches(0.7), Inches(1.8), Inches(6), Inches(3.0), [
    "성형물/벌크 코드·명",
    "투입량 (소수점 셋째자리)",
    "필요 수량 · 이론 필요량",
    "평균 로스율 (5건 가중평균)",
    "최적 제조량",
    "신뢰도 배지 + 점수 + 툴팁",
], size=12, line_spacing=1.3)

add_textbox(s, Inches(7.0), Inches(1.4), Inches(6), Inches(0.4),
            "정렬 & 시각 강조", size=14, bold=True, color=RED)
add_bullet(s, Inches(7.0), Inches(1.8), Inches(6), Inches(3.0), [
    "정렬: 신뢰도 / 제조량 / 로스율 순",
    "저신뢰도 행: 주황 / 빨강 배경",
    "ML vs SAP 차이: 5%p 노랑, 10%p 빨강",
    "5건 가중평균 툴팁 (계산 과정 표시)",
    "이력 모달 75% 확대 (5건 테두리)",
    "행 더블클릭 → 과거 이력 조회",
], size=12, line_spacing=1.3)

# 엑셀 다운로드 컬럼
add_textbox(s, Inches(0.7), Inches(4.85), Inches(12), Inches(0.4),
            "엑셀 다운로드 — 14개 컬럼", size=14, bold=True, color=RED)
add_box(s, Inches(0.7), Inches(5.3), Inches(12.1), Inches(1.6), fill=BG_LIGHT)
add_textbox(s, Inches(0.9), Inches(5.4), Inches(11.7), Inches(0.35),
            "기본 9개 컬럼", size=11, bold=True, color=MID)
add_textbox(s, Inches(0.9), Inches(5.7), Inches(11.7), Inches(0.35),
            "성형물코드 · 성형물명 · 벌크코드 · 벌크명 · 투입량(g) · 필요수량(ea) · 이론 필요량(g) · 평균 로스율 · 최적 제조량(g)",
            size=11, color=DARK)
add_textbox(s, Inches(0.9), Inches(6.15), Inches(11.7), Inches(0.35),
            "신뢰도 관련 5개 컬럼 (신규)", size=11, bold=True, color=RED)
add_textbox(s, Inches(0.9), Inches(6.45), Inches(11.7), Inches(0.35),
            "신뢰도 등급 · 신뢰도 점수 · 유효 표본수 · 편차(%) · 최신 이력(일전)",
            size=11, color=DARK)

slide_footer(s, 9, TOTAL)

# ===================== SLIDE 10: 시스템 특징 & 요약 =====================
s = prs.slides.add_slide(blank)
slide_header(s, 9, "시스템 특징 & 요약")

# 특징 3개
fy = Inches(1.4)
fh = Inches(1.8)
for i, (title, desc, color) in enumerate([
    ("데이터 자동 저장",
     "IndexedDB에 SAP·환입/폐기 데이터를 자동 저장\n페이지 재접속 시 자동 복원되어 매번 재업로드 불필요",
     RED),
    ("운영 유연성",
     "BOM만 초기화 버튼 / 전체 초기화 버튼 분리\n기존 이력 유지한 채 새로운 발주 건만 계산 가능",
     BLUE),
    ("지속적 학습 반영",
     "매주 최신 SAP 실적 데이터 반영\nretrain.py 실행 시 ML 모델 자동 재학습 + 지표 갱신",
     RED),
]):
    x = Inches(0.5 + i * 4.15)
    add_box(s, x, fy, Inches(4.0), fh, fill=BG_LIGHT, line=color)
    add_textbox(s, x, fy+Inches(0.2), Inches(4.0), Inches(0.4), title,
                size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(s, x+Inches(0.2), fy+Inches(0.7), Inches(3.6), Inches(1.0), desc,
                size=11, color=DARK, align=PP_ALIGN.CENTER)

# 향후 발전 방향
add_textbox(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.4),
            "향후 발전 방향", size=14, bold=True, color=RED)
add_bullet(s, Inches(0.7), Inches(3.95), Inches(12), Inches(1.8), [
    ("데이터 축적", "매주 SAP 실적 업데이트로 신뢰도 '매우 낮음' 품목 감소 → 예측 정확도 상승"),
    ("ML 정확도 개선", "현재 테스트 MAE 11.91% → 추가 피처 엔지니어링으로 한 자리 수 목표"),
    ("자동화 확장", "SAP_BOM조회.vbs 통한 발주량 자동 연동, 결과 리포트 자동 생성"),
], size=12, line_spacing=1.3)

# 마무리 박스
add_box(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.1), fill=RGBColor(0xFF, 0xF5, 0xF5))
add_textbox(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.4),
            "요약", size=13, bold=True, color=RED)
add_textbox(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.7),
            "SAP 실적 + ML + 신뢰도 평가를 결합하여, 현장에서 실측 없이 최적 제조량을 빠르게 도출.\n데이터가 쌓일수록 예측 신뢰도가 자동으로 상승하는 구조.",
            size=12, color=DARK)

slide_footer(s, 10, TOTAL)

# 저장
output_path = "/home/djlee/cosmax/project2/파우더_벌크_제조량예측_프로그램_소개.pptx"
prs.save(output_path)
size_kb = os.path.getsize(output_path) / 1024
print(f"저장 완료: {output_path}")
print(f"파일 크기: {size_kb:.1f} KB")
print(f"슬라이드 수: {len(prs.slides)}")
